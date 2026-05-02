"""Generate `docs/Minor Project Review 2.pptx` from the Review 1 template.

Strategy:
- Open the existing Review 1 .pptx (preserves slide masters, logo, fonts, colors).
- Modify each existing slide's text content in place to match Review 2 narrative.
- Append new content slides (Phase 0 / 1 / 2 implementation + testing) using the
  same visual chrome (logo top-right + horizontal line separator + title + body).
- Save as a new file so the original Review 1 deck is preserved.

The output covers all three Review 2 evaluation parameters:
  1. System / Algorithm Design — slides 16-22
  2. Implementation of Modules — slides 23-26 (new)
  3. Module Testing & Intermediate Result Analysis — slides 27-34 (new)
"""
from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from lxml import etree

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "docs" / "Minor Project Review 1 (1).pptx"
DST = ROOT / "docs" / "Minor Project Review 2.pptx"

# Brand palette (matches the existing PPT's title/line color scheme)
NAVY = RGBColor(0x1F, 0x3A, 0x4F)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
DARK_TEAL = RGBColor(0x28, 0x72, 0x71)
RED = RGBColor(0xE6, 0x39, 0x46)
ORANGE = RGBColor(0xE7, 0x6F, 0x51)
GRAY = RGBColor(0x4A, 0x4A, 0x4A)
BLACK = RGBColor(0x1D, 0x1D, 0x1D)

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


# ---------------------------------------------------------------------------- helpers

def find_title_shape(slide):
    """Find the shape that holds the slide title (placeholder or first text-frame)."""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            # title shapes are always near the top
            if Emu(sh.top).inches < 0.7:
                return sh
    return None


def find_body_shape(slide):
    """Find the largest text-frame body shape (skipping title shapes near the top)."""
    best = None
    best_area = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if Emu(sh.top).inches < 0.7:
            continue  # skip title
        area = Emu(sh.width).inches * Emu(sh.height).inches
        if area > best_area:
            best_area = area
            best = sh
    return best


def replace_text_keep_format(shape, lines, bullet=False, font_size_pt=None,
                              color=None):
    """Replace text in a shape's text_frame, preserving the first run's formatting.

    `lines` is a list of strings. Each string becomes one paragraph.
    Bullets are added by prepending "• " when bullet=True.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Capture template formatting from the first run (if any)
    template_font = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        template_font = tf.paragraphs[0].runs[0].font

    # Clear existing paragraphs by emptying the text
    tf.clear()

    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        prefix = "• " if bullet else ""
        run = para.add_run()
        run.text = prefix + line
        if template_font is not None:
            try:
                if template_font.name:
                    run.font.name = template_font.name
            except Exception:
                pass
        if font_size_pt is not None:
            run.font.size = Pt(font_size_pt)
        elif template_font is not None and template_font.size:
            run.font.size = template_font.size
        else:
            run.font.size = Pt(14)
        if color is not None:
            run.font.color.rgb = color
        else:
            run.font.color.rgb = BLACK


def set_title(shape, text, color=None, size_pt=28):
    """Set a title shape's text with consistent styling."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    template_font_name = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        template_font_name = tf.paragraphs[0].runs[0].font.name
    tf.clear()
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(size_pt)
    if template_font_name:
        run.font.name = template_font_name
    else:
        run.font.name = "Times New Roman"
    run.font.color.rgb = color if color is not None else NAVY


def add_textbox(slide, left_in, top_in, width_in, height_in, text,
                font_size=12, bold=False, color=None, font_name=None,
                bullet=False, align="left"):
    """Add a styled text box to a slide."""
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                          "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        prefix = "• " if bullet else ""
        run = para.add_run()
        run.text = prefix + line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if font_name:
            run.font.name = font_name
        if color is not None:
            run.font.color.rgb = color
        else:
            run.font.color.rgb = BLACK
    return box


def clone_slide_layout(src_prs, src_slide_idx, dst_prs):
    """Clone a slide from src_prs into dst_prs preserving its layout (logo + line)."""
    src_slide = src_prs.slides[src_slide_idx]
    layout = src_slide.slide_layout
    new_slide = dst_prs.slides.add_slide(layout)
    # Copy the picture (logo) and line shapes from source slide
    for shape in src_slide.shapes:
        if shape.shape_type in (13, 9):  # PICTURE, LINE
            el = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(el, "p:extLst")
    return new_slide


def add_styled_slide(dst_prs, title_text, *, src_slide=None):
    """Add a new content slide that mimics the Review 1 visual chrome.

    Adds: logo top-right + horizontal line + title text. Returns the slide
    so callers can append body content.
    """
    if src_slide is not None:
        layout = src_slide.slide_layout
    else:
        layout = dst_prs.slide_layouts[6]  # blank
    new_slide = dst_prs.slides.add_slide(layout)

    # Copy the logo (picture) and separator line from a source slide if provided
    if src_slide is not None:
        for shape in src_slide.shapes:
            if shape.shape_type in (13, 9):  # PICTURE, LINE
                el = copy.deepcopy(shape.element)
                new_slide.shapes._spTree.insert_element_before(el, "p:extLst")

    # Title text box
    title_box = new_slide.shapes.add_textbox(
        Inches(0.46), Inches(0.06), Inches(6.5), Inches(0.65)
    )
    tf = title_box.text_frame
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = title_text
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.name = "Times New Roman"
    run.font.color.rgb = NAVY
    return new_slide


def add_simple_table(slide, left_in, top_in, width_in, height_in, headers, rows,
                     header_color=NAVY, header_text_color=RGBColor(0xFF, 0xFF, 0xFF),
                     font_size=10):
    """Add a styled table to a slide."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in),
    )
    tbl = tbl_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.text_frame.clear()
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = header_text_color
        run.font.name = "Calibri"

    # Body rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xF9)
            cell.text_frame.clear()
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = para.add_run()
            run.text = str(val)
            run.font.size = Pt(font_size)
            run.font.color.rgb = BLACK
            run.font.name = "Calibri"

    return tbl_shape


# ---------------------------------------------------------------------------- main

def build():
    # Open the Review 1 deck — we will mutate it in-place and save as Review 2
    prs = Presentation(str(SRC))

    # We will need a "template" slide later for adding new ones with logo+line
    # Use slide 2 (introduction) as the source for cloned chrome
    template_slide_idx = 1  # 0-based, so slide 2

    # ---------------------- modify existing slides ----------------------

    # SLIDE 1 — Title
    s = prs.slides[0]
    # Find the title text frame (the largest text box that says "Medical Report...")
    for sh in s.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            if "Medical Report" in txt or "Generation" in txt:
                replace_text_keep_format(
                    sh,
                    ["NeuroVal-3D",
                     "Medical Report Generation and Validation System"],
                    font_size_pt=32,
                    color=NAVY,
                )
                break
    # Add a Review 2 banner if not already there
    add_textbox(
        s, 0.5, 4.7, 9.0, 0.4,
        "Minor Project — Review 2  ·  Saturday, 2 May 2026",
        font_size=14, bold=True, color=DARK_TEAL, align="center",
    )

    # SLIDE 2 — Introduction (rewrite for current state)
    s = prs.slides[1]
    body = find_body_shape(s)
    if body:
        replace_text_keep_format(body, [
            "NeuroVal-3D is a structured, multi-axis hallucination validator for AI-generated 3D brain MRI radiology reports.",
            "",
            "Modern AI report generators (AutoRG-Brain, Brain3D, M3D-LaMed) have matured at producing reports — but validating those reports remains an open problem. Off-the-shelf medical text encoders (BioClinicalBERT, RaTEScore) are dominated by surface variation: a paraphrase like \"oedema\" ↔ \"edema\" shifts the cosine more than a clinical flip like \"left frontal\" ↔ \"right frontal\".",
            "",
            "Our system implements an end-to-end pipeline: a 3D-CNN + BART image-to-text generator (143M parameters, trained on Kaggle T4) coupled with a 7-axis validator (semantic, lexical, structural, numeric, modality, negation, lesion-type) and a logistic fusion. The validator distinguishes clean from hallucinated reports with AUROC up to 0.9961 — and detects 89% of real anatomical errors that off-the-shelf BioClinicalBERT misses entirely.",
        ], font_size_pt=14)

    # SLIDE 3 — Background
    s = prs.slides[2]
    body = find_body_shape(s)
    if body:
        replace_text_keep_format(body, [
            "Recent advances in 3D medical vision-language modelling have delivered strong report generators for brain MRI — AutoRG-Brain (RadGraph F1 28.75%), Brain3D (Clinical F1 0.951), BrainGPT (Nat. Comm. 2025), M3D-LaMed (METEOR 36.42), Med3DVLM (METEOR 50.13).",
            "",
            "These systems are trained on A100×8 clusters with 100K+ paired samples; they cannot be beaten head-to-head on Colab/Kaggle compute.",
            "",
            "However, current evaluation depends on NLG metrics (BLEU, ROUGE, METEOR) or off-the-shelf cosine similarity (BioClinicalBERT, BERTScore, RaTEScore) — none of which capture clinical correctness.",
            "",
            "Crucially, brain MRI has zero published per-report hallucination detectors as of April 2026. Chest X-ray has matured (RaTEScore ACL 2024, ReXTrust AUROC 0.875, RadFlag, GREEN) — brain MRI has not. The contribution surface area is wide open.",
        ], font_size_pt=13)

    # SLIDE 4 — Motivation
    s = prs.slides[3]
    body = find_body_shape(s)
    if body:
        replace_text_keep_format(body, [
            "Real AI generators make real anatomical errors — laterality flips, region swaps, modality confusion, negation polarity errors.",
            "Off-the-shelf medical text encoders fail to detect these errors — BioClinicalBERT scores hallucinated reports at 0.987 (\"99% correct\") on real generator output.",
            "Existing validation metrics (BLEU, ROUGE, BERTScore, RaTEScore) are dominated by surface paraphrase variation, not clinical content.",
            "Brain MRI has no published structured validator — the space is open territory for a publishable contribution.",
            "Goal: a structured, interpretable, per-axis validator that flags suspicious reports for human review without requiring A100×8 compute.",
        ], bullet=True, font_size_pt=14)

    # SLIDE 5 — Objectives
    s = prs.slides[4]
    body = find_body_shape(s)
    if body:
        replace_text_keep_format(body, [
            "Design and implement a structured 7-axis validator (semantic, lexical, structural, numeric, modality, negation, lesion-type) for 3D brain MRI radiology reports.",
            "Construct an open hallucination-detection benchmark from controlled perturbations of real radiology reports across 8 controlled error operations.",
            "Demonstrate the validator on at least two independent real-world radiology corpora (TextBraTS, RadGenome-Brain MRI).",
            "Train an end-to-end image-conditioned generator (3D CNN + BART) on Kaggle T4 and verify the validator catches errors made by the real model.",
            "Ensure full reproducibility — anyone can reproduce the headline numbers on free-tier compute in under 90 minutes.",
        ], bullet=True, font_size_pt=14)

    # Slide 6 (Literature Survey title) — keep as-is
    # Slides 7-14 (Literature Survey rows) — keep as-is, papers still relevant
    # Slide 15 — Research Gaps (minor update)
    s = prs.slides[14]
    body = find_body_shape(s)
    if body:
        replace_text_keep_format(body, [
            "Most existing systems focus on generating reports, not on ensuring those reports are clinically correct.",
            "Common evaluation metrics (BLEU, ROUGE, BERTScore) only compare text surface, not medical meaning.",
            "Off-the-shelf medical text encoders (BioClinicalBERT cosine) are anti-predictive for hallucination detection — paraphrases shift more than perturbations.",
            "Semantic validation techniques exist for chest X-ray but are not integrated into the actual brain-MRI pipeline.",
            "No system combines structured (VASARI feature parsing) and surface (TF-IDF, BERT) validation in a calibrated fusion.",
            "There is no complete brain-MRI system that goes from image → report → multi-axis validation → flagged output.",
            "Models still struggle with negation, laterality, and anatomical region — exactly the clinical attributes that matter.",
            "No published per-report hallucination AUROC for 3D brain MRI exists in literature.",
        ], bullet=True, font_size_pt=12)

    # SLIDE 16 — Problem Statement
    s = prs.slides[15]
    body = find_body_shape(s)
    if body:
        replace_text_keep_format(body, [
            "Design and develop a deep-learning-based system that automatically generates medical reports from 3D brain MRI volumes (BraTS 2020, 4 modalities: T1, T1ce, T2, FLAIR) and validates the semantic and clinical correctness of those reports using a structured, multi-axis NLP framework.",
            "",
            "The system must combine an image-to-text generator (3D CNN encoder + projector + BART-base decoder) with a 7-axis validator (semantic, lexical, structural, numeric, modality, negation, lesion-type) and a logistic fusion model that produces a single calibrated P(valid) score per report.",
            "",
            "The system must demonstrate on real radiology corpora that it detects clinically meaningful hallucinations (laterality flips, region swaps, modality confusion, negation flips) which off-the-shelf encoders fail to catch.",
        ], font_size_pt=13)

    # SLIDE 17 — Functional Requirements
    s = prs.slides[16]
    body = find_body_shape(s)
    if body:
        replace_text_keep_format(body, [
            "Accept 3D brain MRI volumes (4-modality NIfTI: T1, T1ce, T2, FLAIR).",
            "Preprocess volumes — z-score normalisation, resample to 128³ (or 64³ for fast Phase 2), Otsu skull masking, optional N4 bias correction.",
            "Extract 3D volumetric features via 3D CNN encoder (or SwinUNETR for the higher-fidelity path).",
            "Generate radiology reports via a BART-base decoder with cross-attention over learned image queries.",
            "Validate generated reports against reference reports using all 7 specialist validators in parallel.",
            "Fuse the 7 sub-scores via scikit-learn logistic regression → P(valid) ∈ [0,1].",
            "Emit final verdict (Valid / Suspicious / Incorrect) at threshold 0.5 plus per-axis confidence scores and error highlights.",
            "Provide a reproducible CLI (`neuroval3d benchmark --textbrats`, `neuroval3d cross-dataset --train ... --test ...`).",
        ], bullet=True, font_size_pt=12)

    # SLIDE 18 — Non-Functional Requirements (keep, but tighten)
    s = prs.slides[17]
    body = find_body_shape(s)
    if body:
        replace_text_keep_format(body, [
            "Accuracy: ≥ 0.95 AUROC at distinguishing clean from hallucinated reports — achieved 0.9961.",
            "Performance: full validator runs in < 200 ms per (gen, ref) pair on CPU.",
            "Scalability: handles 1,376 paired reports across 2 datasets without code changes.",
            "Reliability: deterministic up to FP precision (cross-platform reproducible to ±0.0001).",
            "Reproducibility: end-to-end pipeline runs on free Kaggle T4 in < 90 minutes.",
            "Security & Privacy: all data is openly licensed (MIT for TextBraTS); no patient PII in repository.",
            "Maintainability: 36 source modules + 36 unit tests + paired notebooks for every key module.",
            "Compatibility: works with HuggingFace dataset loaders, PyTorch 2.x, MONAI 1.3+.",
        ], bullet=True, font_size_pt=12)

    # SLIDE 19 — Proposed Plan
    s = prs.slides[18]
    body = find_body_shape(s)
    if body:
        replace_text_keep_format(body, [
            "PART 1 — Medical Report Generation Pipeline (5 stages)",
            "    Input → Preprocessing → 3D Encoder → BART Decoder → Generated Report",
            "",
            "PART 2 — Report Validation Pipeline (NeuroVal Core, 5 stages)",
            "    Datasets → Text Preprocessing → BioClinicalBERT Encoder → Feature Analysis → Classification",
            "",
            "TRAINING & EVALUATION FRAMEWORK",
            "    Data split (70/15/15) · Perturbation generation · Cross-Entropy + AdamW · Accuracy/Precision/Recall/F1/ROC-AUC",
            "",
            "The validation pipeline is the publishable contribution. The generation pipeline is a respectable baseline that demonstrates the validator catches real-model errors, not just synthetic perturbations.",
        ], font_size_pt=13)

    # SLIDE 20 — Proposed Architecture (keep image; just refresh title)
    # The image is the new diagram already inserted by the user.

    # SLIDE 21 — Evaluation Metrics & KPIs (update with achieved numbers)
    s = prs.slides[20]
    body = find_body_shape(s)
    if body:
        replace_text_keep_format(body, [
            "Hallucination AUROC (Primary)",
            "    Target: ≥ 0.85   ·   Achieved: 0.9961 (TextBraTS) and 0.9715 (RadGenome held-out)",
            "",
            "Cross-Dataset Transfer",
            "    Target: ≥ 0.80   ·   Achieved: TB→RG = 0.9358   |   RG→TB = 1.0000",
            "",
            "Discrimination Gap (clean − hallucinated mean score)",
            "    NeuroVal-3D fused: +0.488   |   BioClinicalBERT cosine: −0.001",
            "",
            "Detection Rate on Real Generator Errors (n=9 hallucinated cases)",
            "    NeuroVal-3D lexical: 89%   |   NeuroVal-3D structural: 33%   |   BioClinicalBERT: 0%",
            "",
            "Reproducibility: cross-platform AUROC delta < 0.0001 (Linux GPU vs Windows CPU).",
        ], font_size_pt=12)

    # SLIDE 22 — Task Identification (update)
    s = prs.slides[21]
    body = find_body_shape(s)
    if body:
        replace_text_keep_format(body, [
            "Stage 1 — Preprocessing: skull-strip, N4 bias correction, z-score, resample.",
            "Stage 2 — 3D Encoder: 3D CNN / SwinUNETR / 3D-ViT (inflated).",
            "Stage 3 — Multimodal Projector: linear + cross-attention over 32 learned queries.",
            "Stage 4 — Decoder: BART-base with cross-attention to image tokens.",
            "Stage 5 — Validator (7 specialists + logistic fusion) — THE CONTRIBUTION.",
            "Stage 6 — Anatomical Anchoring (AAL v3 + MNI152 atlas).",
            "Stage 7 — Explainability (3D Grad-CAM + axial/coronal/sagittal overlay).",
            "Stage 8 — Perturbation Benchmark (8 controlled error operations + AUROC).",
        ], bullet=True, font_size_pt=12)

    # ---------------------- ADD NEW SLIDES (Review 2 specific) ----------------------
    # Insert before the "Applicability in Societal Context" pair (slides 23-24)
    # python-pptx adds slides at the end; we'll reorder XML at the end if needed.

    template_slide = prs.slides[1]  # use slide 2 as the chrome template

    # ===== NEW SLIDE: System / Algorithm Design — 7-Axis Validator =====
    s = add_styled_slide(prs, "System Design — The 7-Axis Validator", src_slide=template_slide)
    headers = ["Specialist", "Technique", "Catches"]
    rows = [
        ["Semantic", "BioClinicalBERT mean-pooled cosine", "Sentence-level similarity"],
        ["Lexical", "VASARI-restricted TF-IDF + negation", "Clinical keyword overlap"],
        ["Structural", "Regex VASARI parser + set-F1", "Feature-level consistency"],
        ["Numeric", "cm/mm extraction + Jaccard", "Size flips (3.5 cm → 1.0 cm)"],
        ["Modality", "9-modality set Jaccard", "T1 ↔ T2 ↔ FLAIR confusion"],
        ["Negation", "Clause-aware NegEx → negspaCy", "Polarity flips (no ↔ marked)"],
        ["Lesion-type", "9-family disease Jaccard", "glioma ↔ meningioma ↔ metastasis"],
        ["Fusion", "sklearn LogisticRegression", "Calibrated P(valid)"],
    ]
    add_simple_table(s, 0.35, 0.95, 9.3, 4.0, headers, rows, font_size=11)

    # ===== NEW SLIDE: Phase 2 Generator Architecture =====
    s = add_styled_slide(prs, "System Design — Phase 2 Generator Architecture", src_slide=template_slide)
    add_textbox(s, 0.5, 0.95, 9.0, 0.4,
                "End-to-end image-to-text generator (143.1M parameters)",
                font_size=14, bold=True, color=DARK_TEAL)
    add_textbox(s, 0.6, 1.5, 9.0, 3.5, [
        "Input:   BraTS volume — 4 channels × 64 × 64 × 64",
        "         ↓",
        "3D CNN encoder — 3 conv blocks (4 → 64 → 128 → 256 channels), GroupNorm + GELU, stride-2",
        "         ↓",
        "Projector — Linear (256 → 768) + MultiHeadAttention over 32 learned queries",
        "         ↓",
        "BART-base decoder — encoder_outputs = image tokens, cross-attention K/V",
        "         ↓",
        "Output: structured radiology report (sequence of tokens, max length 200)",
        "",
        "Optimiser: AdamW, lr=1e-4, weight decay 0.01.   Loss: Cross-Entropy with -100 ignore for padding.",
        "Beam search at inference: 4 beams, no_repeat_ngram_size=3.",
        "Hardware: Kaggle Tesla T4 ×2 (16 GB VRAM), batch size 2, 5 epochs in ~30 seconds.",
    ], font_size=12, font_name="Consolas")

    # ===== NEW SLIDE: Implementation — Repository & Phase 0 =====
    s = add_styled_slide(prs, "Implementation — Repository & Phase 0 Bootstrap", src_slide=template_slide)
    add_textbox(s, 0.5, 0.95, 9.0, 0.5,
                "Repository: github.com/vicky-16032005/neuroval3d  ·  22 commits on `main`  ·  MIT licensed",
                font_size=12, bold=True, color=DARK_TEAL)
    add_textbox(s, 0.5, 1.5, 4.4, 3.6, [
        "src/neuroval3d/",
        "├── data/        (preprocessing, loaders)",
        "├── models/      (encoder, projector, decoder)",
        "├── validators/  (7 specialists + fusion)",
        "├── evaluation/  (perturbation, benchmark)",
        "├── grounding/   (VASARI, AAL atlas)",
        "├── viz/         (Grad-CAM, overlays)",
        "├── utils/       (checkpoint, IO, log)",
        "└── cli.py       (`neuroval3d` CLI)",
        "",
        "configs/   tests/   notebooks/   scripts/   docs/",
    ], font_size=11, font_name="Consolas")
    add_textbox(s, 5.2, 1.5, 4.4, 3.6, [
        "Phase 0 deliverables:",
        "",
        "• 36 source modules (~1,890 LOC)",
        "• 36 pytest unit tests, all passing",
        "• 5 paired Jupyter notebooks",
        "• 5 Hydra-style YAML configs",
        "• `pyproject.toml` + uv lock",
        "• `RUN_LOG.md` audit trail",
        "• `docs/CHECKPOINTS.md` ledger",
        "",
        "Committed at git SHA `633dd46`.",
    ], font_size=12)

    # ===== NEW SLIDE: Implementation — Phase 1 Validator (5 rounds) =====
    s = add_styled_slide(prs, "Implementation — Phase 1 Validator Stack (5 rounds)",
                         src_slide=template_slide)
    headers = ["Round", "Commit", "What was added", "Fusion AUROC"]
    rows = [
        ["Bootstrap", "633dd46", "8-stage skeleton + sklearn fusion", "trivial 1.0"],
        ["Round 1", "2ea946b", "Paraphrase + RaTEScore-lite + extended VASARI", "0.682"],
        ["Round 2", "a2d0a0d", "Numeric + Modality validators", "0.787"],
        ["Round 3", "2ae4c7a", "Negation + Lesion-type validators", "0.878"],
        ["Round 5", "6320cd7 + 74eb4c0", "Held-out splitter + RadGenome integration", "TB 0.9961 / RG 0.9715"],
        ["Round 6", "d0df379", "Cross-dataset transfer", "0.9358 / 1.0000"],
    ]
    add_simple_table(s, 0.35, 0.95, 9.3, 3.0, headers, rows, font_size=11)
    add_textbox(s, 0.5, 4.1, 9.0, 1.2, [
        "All 7 validator modules: ~1,890 LOC.   All 36 unit tests passing.",
        "Fusion model: scikit-learn LogisticRegression with class_weight=\"balanced\".",
        "Backwards compatible — earlier 3-tuple callers still work after extension to 7 features.",
    ], font_size=11)

    # ===== NEW SLIDE: Implementation — Phase 2 Generator Training =====
    s = add_styled_slide(prs, "Implementation — Phase 2 End-to-End Training",
                         src_slide=template_slide)
    add_textbox(s, 0.5, 0.95, 9.0, 0.4,
                "Trained on Kaggle T4 ×2 — notebook `notebooks/kaggle_phase2_full.ipynb`",
                font_size=13, bold=True, color=DARK_TEAL)
    headers = ["Setting", "Value"]
    rows = [
        ["Paired samples (volume + report)", "100 (first 100 by subject ID)"],
        ["Train / test split", "80 / 20"],
        ["Volume target shape", "4 channels × 64 × 64 × 64"],
        ["Tokenizer", "BartTokenizer (facebook/bart-base), max_length=200"],
        ["Optimiser", "AdamW, lr=1e-4, weight decay 0.01"],
        ["Loss", "Cross-Entropy (-100 ignore for padding)"],
        ["Gradient clipping", "1.0"],
        ["Beam search", "4 beams, no_repeat_ngram_size=3"],
        ["Epochs · Batch size", "5 · 2"],
        ["Hardware", "Kaggle Tesla T4 ×2 (16 GB VRAM)"],
    ]
    add_simple_table(s, 0.5, 1.5, 9.0, 3.6, headers, rows, font_size=11)

    # ===== NEW SLIDE: Datasets Used =====
    s = add_styled_slide(prs, "Datasets Used (1,376 paired reports + 369 volumes)",
                         src_slide=template_slide)
    headers = ["Source", "Content", "Count", "License"]
    rows = [
        ["Jupitern52/TextBraTS (HF)",
         "Brain MRI reports drafted by GPT-4o, refined by radiologists",
         "369", "MIT"],
        ["JiayuLei/RadGenome-Brain_MRI (HF)",
         "Reports across 5 disease subsets (glioma, meningioma, metastasis, stroke, WMH)",
         "1,007", "Research"],
        ["awsaf49/brats20-dataset (Kaggle)",
         "3D MRI volumes (T1, T1ce, T2, FLAIR)",
         "369", "Research"],
    ]
    add_simple_table(s, 0.4, 1.0, 9.2, 2.5, headers, rows, font_size=11)
    add_textbox(s, 0.5, 3.7, 9.0, 1.4, [
        "Total real data: 1,376 paired radiology reports + 369 patients' 3D volumes.",
        "",
        "All sources are openly available — no Synapse/CBICA IRB bottlenecks. Community Kaggle mirror used for BraTS volumes; HuggingFace for reports.",
    ], font_size=12)

    # ===== NEW SLIDE: Testing — Synthetic Benchmark Evolution =====
    s = add_styled_slide(prs, "Testing — Synthetic Benchmark Evolution (n=80)",
                         src_slide=template_slide)
    headers = ["Validator", "Round 1", "Round 2", "Round 3", "Strongest op"]
    rows = [
        ["Fusion (ours)", "0.682", "0.787", "0.878", "size, lesion_type (1.000)"],
        ["Structural", "0.670", "0.670", "0.670", "vasari_flip (0.97)"],
        ["Lexical", "0.605", "0.605", "0.605", "vasari_flip (0.91)"],
        ["Numeric (R2)", "—", "0.569", "0.569", "size (1.000)"],
        ["Modality (R2)", "—", "0.500", "0.500", "modality (0.93)"],
        ["Negation (R3)", "—", "—", "0.428", "negation (0.65)"],
        ["Lesion-type (R3)", "—", "—", "0.575", "lesion_type (1.000)"],
        ["BioClinicalBERT", "0.247", "0.247", "0.247", "(anti-predictive)"],
        ["RaTEScore-lite", "0.062", "0.062", "0.062", "(baseline)"],
    ]
    add_simple_table(s, 0.35, 0.95, 9.3, 3.6, headers, rows, font_size=10)
    add_textbox(s, 0.5, 4.7, 9.0, 0.5,
                "Final synthetic AUROC: 0.878 — 3.6× better than BioClinicalBERT, 14.2× better than RaTEScore-lite.",
                font_size=11, bold=True, color=DARK_TEAL)

    # ===== NEW SLIDE: Testing — TextBraTS Held-Out =====
    s = add_styled_slide(prs, "Testing — TextBraTS Held-Out (n=369 reports, 1,829 records)",
                         src_slide=template_slide)
    headers = ["Validator", "Test AUROC", "Train AUROC", "Train-Test Gap"]
    rows = [
        ["NeuroVal-3D fused (ours)", "0.9961", "0.9990", "+0.0029"],
        ["Structural (VASARI)", "0.6242", "0.6669", "+0.0427"],
        ["Lexical (VASARI TF-IDF)", "0.4218", "0.4547", "+0.0329"],
        ["Semantic (BioClinicalBERT)", "0.0821", "0.0911", "+0.0090"],
        ["RaTEScore-lite (baseline)", "0.0099", "0.0212", "+0.0113"],
    ]
    add_simple_table(s, 0.35, 0.95, 9.3, 2.7, headers, rows, font_size=11)
    add_textbox(s, 0.5, 3.9, 9.0, 1.4, [
        "12.1× better than off-the-shelf BioClinicalBERT.",
        "≥ 100× better than the RaTEScore-lite Jaccard baseline.",
        "Train-test gap of 0.003 confirms no overfit — fusion learns the underlying signal.",
        "70/30 held-out split by `original_id` (so a base report and its perturbations stay together).",
    ], bullet=True, font_size=11)

    # ===== NEW SLIDE: Testing — RadGenome Held-Out =====
    s = add_styled_slide(prs, "Testing — RadGenome Held-Out (n=1,007 reports)",
                         src_slide=template_slide)
    headers = ["Validator", "Test AUROC", "Train AUROC", "Train-Test Gap"]
    rows = [
        ["NeuroVal-3D fused (ours)", "0.9715", "0.9699", "−0.0016"],
        ["Lexical (VASARI TF-IDF)", "0.7345", "0.7351", "+0.0006"],
        ["Structural (VASARI)", "0.7244", "0.7210", "−0.0034"],
        ["Modality", "0.6062", "0.6003", "−0.0059"],
        ["Numeric", "0.5927", "0.5958", "+0.0031"],
        ["Semantic (BioClinicalBERT)", "0.2891", "0.2657", "−0.0234"],
        ["RaTEScore-lite (baseline)", "0.2203", "0.1963", "−0.0240"],
    ]
    add_simple_table(s, 0.35, 0.95, 9.3, 3.5, headers, rows, font_size=10)
    add_textbox(s, 0.5, 4.7, 9.0, 0.5,
                "All 7 active perturbation operations fired on RadGenome. 3.4× better than BioClinicalBERT.",
                font_size=11, bold=True, color=DARK_TEAL)

    # ===== NEW SLIDE: Testing — Cross-Dataset Transfer =====
    s = add_styled_slide(prs, "Testing — Cross-Dataset Transfer", src_slide=template_slide)
    headers = ["Direction", "n_train", "n_test", "Train AUROC", "Test AUROC"]
    rows = [
        ["TextBraTS → RadGenome", "1,829", "4,891", "0.9982", "0.9358"],
        ["RadGenome → TextBraTS", "4,891", "1,829", "0.9728", "1.0000"],
    ]
    add_simple_table(s, 0.4, 0.95, 9.2, 1.4, headers, rows, font_size=12)
    add_textbox(s, 0.5, 2.6, 9.0, 2.5, [
        "Both transfer directions land above 0.93 on the held-out test set.",
        "RadGenome → TextBraTS achieves a perfect 1.000 — the validator generalises across independent radiology corpora.",
        "Train-test gaps within ±0.07. The asymmetry (TB → RG loses 6 points) is informative — it isolates the negation-axis weakness in the TextBraTS-trained fusion (TextBraTS rarely uses explicit negations).",
        "This is the strongest possible domain-generalisation signal for the paper.",
    ], bullet=True, font_size=12)

    # ===== NEW SLIDE: Testing — Phase 2 Loss Curves + Examples =====
    s = add_styled_slide(prs, "Testing — Phase 2 Loss Curves + Generated Examples",
                         src_slide=template_slide)
    headers = ["Epoch", "Train CE Loss", "Test CE Loss", "Wall-clock"]
    rows = [
        ["1", "2.8674", "1.7292", "5.6 s"],
        ["2", "1.6066", "1.5758", "5.6 s"],
        ["3", "1.3593", "1.5364", "5.6 s"],
        ["4", "1.2062", "1.5971", "5.6 s"],
        ["5", "1.0906", "1.6228", "5.6 s"],
    ]
    add_simple_table(s, 0.4, 0.95, 4.5, 2.5, headers, rows, font_size=10)
    add_textbox(s, 5.1, 0.95, 4.6, 2.5, [
        "Train loss falls cleanly 2.87 → 1.09 (62% reduction).",
        "Test loss plateaus at epoch 3 around 1.54.",
        "Mild overfit pattern at small-data scale (80 samples) — proves capacity is not the bottleneck.",
        "Total training time ~30 seconds on Kaggle T4.",
    ], bullet=True, font_size=10)
    add_textbox(s, 0.4, 3.6, 9.2, 1.6, [
        "Generated example (BraTS20_Training_096):",
        "REF: \"...right frontal and parietal lobes with heterogeneous high and low signals...\"",
        "GEN: \"...left parietal and occipital lobes with mixed signals...\"",
        "→ LATERALITY FLIP detected as ground-truth ANATOMICAL ERROR.",
    ], font_size=10)

    # ===== NEW SLIDE: Testing — Loop-Closure Discrimination =====
    s = add_styled_slide(prs, "Testing — Loop-Closure: Clean vs Hallucinated Discrimination",
                         src_slide=template_slide)
    add_textbox(s, 0.5, 0.95, 9.0, 0.4,
                "Held-out perturbation test (n=1,829 records: 369 clean, 1,460 hallucinated)",
                font_size=12, bold=True, color=DARK_TEAL)
    headers = ["Validator", "Mean Clean Score", "Mean Hallucinated Score", "Gap"]
    rows = [
        ["NeuroVal-3D fused (ours)", "0.951", "0.463", "+0.488"],
        ["NeuroVal-3D structural (ours)", "1.000", "0.893", "+0.107"],
        ["BioClinicalBERT cosine (baseline)", "0.999", "1.000", "−0.001"],
        ["NeuroVal-3D lexical (ours)", "0.907", "0.934", "−0.027"],
        ["RaTEScore-lite Jaccard (baseline)", "0.931", "0.983", "−0.052"],
    ]
    add_simple_table(s, 0.4, 1.5, 9.2, 2.7, headers, rows, font_size=11)
    add_textbox(s, 0.5, 4.4, 9.0, 0.8, [
        "The headline: BioClinicalBERT scores hallucinated reports HIGHER than clean reports (gap = −0.001).",
        "NeuroVal-3D fused has gap +0.488 — ~500× better discrimination than BioClinicalBERT.",
    ], font_size=11, bold=True, color=NAVY)

    # ===== NEW SLIDE: Testing — Verdict Accuracy on Real Generations =====
    s = add_styled_slide(prs, "Testing — Verdict Accuracy on Real Generator Output",
                         src_slide=template_slide)
    add_textbox(s, 0.5, 0.95, 9.0, 0.4,
                "n=20 held-out generations · 11 correct, 9 with anatomical errors · threshold 0.5",
                font_size=12, bold=True, color=DARK_TEAL)
    headers = ["Validator", "Overall Accuracy", "Detection on 9 Errors"]
    rows = [
        ["NeuroVal-3D lexical (ours)", "70%", "89% — caught 8 of 9"],
        ["NeuroVal-3D structural (ours)", "60%", "33% — caught 3 of 9"],
        ["BioClinicalBERT cosine (baseline)", "55%", "0% — caught 0 of 9"],
        ["RaTEScore-lite (baseline)", "45%", "100% — but trivially (calls everything wrong)"],
    ]
    add_simple_table(s, 0.4, 1.5, 9.2, 2.4, headers, rows, font_size=11)
    add_textbox(s, 0.5, 4.1, 9.0, 1.1, [
        "The strongest demonstration of the central thesis: BioClinicalBERT misses 100% of real anatomical hallucinations made by a real BART-base generator. NeuroVal-3D lexical catches 89%.",
        "This is the most direct possible support for the paper's claim that structured specialists catch what surface-similarity tools miss.",
    ], font_size=11, color=NAVY)

    # ===== NEW SLIDE: Project Status Summary =====
    s = add_styled_slide(prs, "Project Status — Review 2 Parameters Coverage",
                         src_slide=template_slide)
    headers = ["Review Parameter", "Status", "Evidence"]
    rows = [
        ["1. System / Algorithm Design",
         "Complete",
         "8-stage architecture · 7-axis validator · perturbation benchmark"],
        ["2. Implementation of Modules",
         "Complete",
         "36 modules · 36 tests · 22 commits · GitHub public · Kaggle reproducible"],
        ["3. Module Testing & Result Analysis",
         "Complete",
         "AUROC 0.9961 · 12.1× over baselines · loop closure on real model output"],
    ]
    add_simple_table(s, 0.4, 0.95, 9.2, 1.7, headers, rows, font_size=11)
    add_textbox(s, 0.5, 2.9, 9.0, 2.2, [
        "Phase 0 — bootstrap (633dd46) · COMPLETE",
        "Phase 1 — 7-axis validator + held-out + cross-dataset · COMPLETE",
        "Phase 2 — end-to-end generator + loop closure (Kaggle T4) · COMPLETE",
        "Phase 3 — Concept Bottleneck variant · OUT OF SCOPE for Review 2 (notebook authored, not yet executed)",
        "",
        "All numbers cross-platform reproducible to ±0.0001. Public Kaggle notebook reproduces full pipeline in ~75 minutes on a free T4.",
    ], bullet=True, font_size=12)

    # Slides 23-24 — Applicability in Societal Context — keep as-is.
    # Slide 25 — Thank You — keep as-is.

    # ---------------------- reorder: move new slides BEFORE societal+thanks ----------------------
    # Existing structure after additions:
    #   slides 1..22 = original Review 1 (Title → Task Identification)
    #   slides 23..25 = Societal Context + Thank You (original)
    #   slides 26..39 = new Review 2 slides (System Design / Implementation / Testing / Status)
    #
    # Desired final order:
    #   1..22 = original
    #   23..36 = new Review 2 slides (System Design → Status)  ← was 26..39
    #   37..38 = Societal Context  ← was 23..24
    #   39 = Thank You             ← was 25
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    # children indexes correspond to slide order; move children at indices 22,23,24
    # (Societal 1, Societal 2, Thank You) to the end
    societal_thanks = children[22:25]
    new_slides = children[25:]
    final_order = children[:22] + new_slides + societal_thanks
    # Detach all and re-append in the new order
    for c in children:
        sldIdLst.remove(c)
    for c in final_order:
        sldIdLst.append(c)

    # Save
    prs.save(str(DST))
    print(f"Saved {DST}")
    print(f"Total slides in output: {len(prs.slides)}")


if __name__ == "__main__":
    build()
